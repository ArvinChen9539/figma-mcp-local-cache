from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_cool_presentation(filename="Figma_MCP_Local_Cache_Intro.pptx"):
    prs = Presentation()
    
    # --- Slide 1: Title Slide ---
    # 使用空白布局自定义设计
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    
    # 背景色 - 深空蓝
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47)
    
    # 标题 - 霓虹绿
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "FIGMA MCP LOCAL CACHE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 255, 218)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题 - 亮白
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf_sub = subtitle_box.text_frame
    p_sub = tf_sub.add_paragraph()
    p_sub.text = "极速 · 智能 · 本地化"
    p_sub.font.size = Pt(32)
    p_sub.font.color.rgb = RGBColor(230, 241, 255)
    p_sub.alignment = PP_ALIGN.CENTER

    # --- Slide 2: The Problem (Pain Points) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(23, 42, 69) # 稍浅的深蓝
    
    # 标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = title_shape.text_frame.add_paragraph()
    p.text = "为什么我们需要它？"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 107, 107) # 警告红

    # 痛点列表
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    tf = content_box.text_frame
    
    points = [
        "🐢 API 调用太慢，等待让人抓狂",
        "💸 频繁请求触发限流，业务中断",
        "🚫 断网 = 断粮，无法离线工作",
        "🤯 数据结构复杂，解析耗时耗力"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(136, 146, 176)
        p.space_after = Pt(20)

    # --- Slide 3: The Solution (Core Features) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(10, 25, 47)

    # 标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = title_shape.text_frame.add_paragraph()
    p.text = "核心黑科技"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 255, 218)

    # 功能卡片模拟
    # 卡片1: MCP 协议
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(2), Inches(4), Inches(2.5)) # 矩形
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(17, 34, 64)
    shape.line.color.rgb = RGBColor(100, 255, 218)
    p = shape.text_frame.add_paragraph()
    p.text = "MCP 协议驱动\n\n完美对接 Cursor/Trae\n让 AI 直接读取设计稿"
    p.font.color.rgb = RGBColor(230, 241, 255)
    p.alignment = PP_ALIGN.CENTER

    # 卡片2: 双模存储
    shape = slide.shapes.add_shape(1, Inches(5.5), Inches(2), Inches(4), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(17, 34, 64)
    shape.line.color.rgb = RGBColor(100, 255, 218)
    p = shape.text_frame.add_paragraph()
    p.text = "双模存储引擎\n\nMySQL 数据库 (企业级)\nJSON 文件系统 (轻量级)"
    p.font.color.rgb = RGBColor(230, 241, 255)
    p.alignment = PP_ALIGN.CENTER

    # 卡片3: 智能缓存
    shape = slide.shapes.add_shape(1, Inches(3), Inches(5), Inches(4), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(17, 34, 64)
    shape.line.color.rgb = RGBColor(100, 255, 218)
    p = shape.text_frame.add_paragraph()
    p.text = "智能缓存策略\n\n优先本地命中\n按需强制同步"
    p.font.color.rgb = RGBColor(230, 241, 255)
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 4: Architecture (Visual) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(23, 42, 69)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = title_shape.text_frame.add_paragraph()
    p.text = "系统架构一览"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 206, 107) # 亮黄

    # 简单的架构图绘制
    # AI Agent
    agent_box = slide.shapes.add_shape(1, Inches(1), Inches(3), Inches(2), Inches(1.5))
    agent_box.text = "AI Agent\n(Cursor/Trae)"
    agent_box.fill.solid()  # Fix: Initialize fill type first
    agent_box.fill.fore_color.rgb = RGBColor(100, 255, 218)
    agent_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(10, 25, 47)

    # MCP Server
    mcp_box = slide.shapes.add_shape(1, Inches(4), Inches(3), Inches(2), Inches(1.5))
    mcp_box.text = "Figma MCP\nServer"
    mcp_box.fill.solid() # Fix: Initialize fill type first
    mcp_box.fill.fore_color.rgb = RGBColor(255, 107, 107)
    
    # Storage
    db_box = slide.shapes.add_shape(1, Inches(7), Inches(2), Inches(2), Inches(1.5))
    db_box.text = "MySQL / JSON\nCache"
    db_box.fill.solid() # Fix: Initialize fill type first
    db_box.fill.fore_color.rgb = RGBColor(136, 146, 176)

    # Figma API
    api_box = slide.shapes.add_shape(1, Inches(7), Inches(4.5), Inches(2), Inches(1.5))
    api_box.text = "Figma Cloud\nAPI"
    api_box.fill.solid() # Fix: Initialize fill type first
    api_box.fill.fore_color.rgb = RGBColor(136, 146, 176)

    # Arrows
    arrow = slide.shapes.add_connector(1, Inches(3), Inches(3.75), Inches(4), Inches(3.75)) # Agent -> MCP
    arrow = slide.shapes.add_connector(1, Inches(6), Inches(3.5), Inches(7), Inches(2.75)) # MCP -> DB
    arrow = slide.shapes.add_connector(1, Inches(6), Inches(4), Inches(7), Inches(5.25)) # MCP -> API

    # --- Slide 5: Call to Action ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(10, 25, 47)

    center_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tf = center_box.text_frame
    p = tf.add_paragraph()
    p.text = "立即体验极速开发"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 255, 218)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "github.com/figma-mcp-local-cache"
    p2.font.size = Pt(30)
    p2.font.color.rgb = RGBColor(230, 241, 255)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(30)

    prs.save(filename)
    print(f"Presentation saved to {filename}")

if __name__ == "__main__":
    create_cool_presentation()
